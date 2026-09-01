import os
import re
import csv
import json
import base64
import uuid
import logging
from typing import List, Dict, Any, Optional
from sqlalchemy.orm import Session

# Setup logging
logger = logging.getLogger("api.services.document_parser")

class DocumentParser:
    @classmethod
    def parse_file(
        cls,
        file_path: str,
        file_type: str,
        db: Session,
        organization_id: uuid.UUID,
        user_id: uuid.UUID,
    ) -> Dict[str, Any]:
        """
        Routing method to parse files depending on type and returns extracted content & metadata.
        """
        import time
        import hashlib
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
            
        ftype = file_type.lower().strip(".")
        start_time = time.perf_counter()
        
        # Calculate SHA256 checksum/hash
        hasher = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hasher.update(chunk)
        checksum = hasher.hexdigest()
        
        page_count = 1
        author = "System"
        
        # Try extracting page count / author metadata from structural files
        if ftype == "pdf":
            try:
                import pypdf
                reader = pypdf.PdfReader(file_path)
                page_count = len(reader.pages)
                if reader.metadata:
                    author = reader.metadata.get("/Author", "System")
            except Exception:
                pass
        elif ftype == "docx":
            try:
                import docx
                doc = docx.Document(file_path)
                author = doc.core_properties.author or "System"
            except Exception:
                pass
                
        try:
            if ftype in ["txt", "text"]:
                raw_text = cls._parse_txt(file_path)
            elif ftype in ["md", "markdown"]:
                raw_text = cls._parse_markdown(file_path)
            elif ftype == "csv":
                raw_text = cls._parse_csv(file_path)
            elif ftype in ["xls", "xlsx"]:
                raw_text = cls._parse_excel(file_path)
            elif ftype == "docx":
                raw_text = cls._parse_docx(file_path)
            elif ftype == "pptx":
                raw_text = cls._parse_pptx(file_path)
            elif ftype == "json":
                raw_text = cls._parse_json(file_path)
            elif ftype == "xml":
                raw_text = cls._parse_xml(file_path)
            elif ftype == "html":
                raw_text = cls._parse_html(file_path)
            elif ftype == "rtf":
                raw_text = cls._parse_rtf(file_path)
            elif ftype in ["png", "jpg", "jpeg", "webp", "tiff", "tif"]:
                raw_text = cls._parse_image_ocr(file_path, db, organization_id, user_id)
            elif ftype == "pdf":
                raw_text = cls._parse_pdf(file_path, db, organization_id, user_id)
            else:
                raw_text = cls._parse_txt(file_path)
        except Exception as e:
            logger.error(f"Error parsing file {file_path} of type {file_type}: {e}")
            raise RuntimeError(f"Document parsing failed: {str(e)}")
            
        duration_ms = int((time.perf_counter() - start_time) * 1000)
        language = cls._detect_language(raw_text)
        keywords = cls._extract_keywords(raw_text)
        
        return {
            "content": raw_text,
            "checksum": checksum,
            "page_count": page_count,
            "author": author,
            "language": language,
            "keywords": keywords,
            "duration_ms": duration_ms,
        }

    @classmethod
    def _detect_language(cls, text: str) -> str:
        if not text:
            return "en"
        # Word frequency list overlap checks
        spanish_words = {"y", "el", "la", "de", "que", "en", "un", "los", "se", "del"}
        french_words = {"et", "le", "la", "de", "que", "en", "un", "les", "se", "des"}
        german_words = {"und", "der", "die", "das", "von", "dass", "in", "ein", "zu", "dem"}
        
        words = [w.lower() for w in text.split()[:200]]
        word_set = set(words)
        
        if word_set.intersection(spanish_words):
            return "es"
        elif word_set.intersection(french_words):
            return "fr"
        elif word_set.intersection(german_words):
            return "de"
        return "en"

    @classmethod
    def _extract_keywords(cls, text: str) -> List[str]:
        if not text:
            return []
        stopwords = {"about", "above", "across", "after", "again", "against", "all", "almost", "alone", "along", "already", "also", "although", "always", "among", "amount", "another", "answer", "any", "anyone", "anything", "anyway", "anywhere", "apply", "area", "around", "back", "became", "because", "become", "becomes", "becoming", "before", "behind", "being", "below", "between", "bill", "both", "bottom", "cannot", "could", "detail", "during", "either", "enough", "etc", "even", "ever", "every", "everyone", "everything", "everywhere", "except", "fill", "find", "fire", "first", "former", "formerly", "found", "from", "front", "full", "further", "give", "here", "hereafter", "hereby", "herein", "hereupon", "hers", "herself", "himself", "however", "interest", "into", "keep", "last", "latter", "latterly", "least", "less", "many", "may", "more", "moreover", "most", "mostly", "move", "much", "must", "myself", "name", "namely", "neither", "never", "nevertheless", "next", "nine", "nobody", "none", "noone", "nothing", "now", "nowhere", "often", "once", "only", "onto", "other", "others", "otherwise", "ours", "ourselves", "part", "people", "perhaps", "please", "public", "rather", "same", "seem", "seemed", "seeming", "seems", "serious", "several", "shall", "should", "show", "side", "since", "sincere", "sixty", "someone", "something", "sometime", "sometimes", "somewhere", "still", "such", "system", "take", "than", "that", "the", "their", "them", "themselves", "then", "thence", "there", "thereafter", "thereby", "therefore", "therein", "thereupon", "these", "they", "thick", "thin", "third", "this", "those", "though", "three", "through", "throughout", "thru", "thus", "together", "too", "toward", "towards", "twelve", "twenty", "very", "via", "was", "we", "well", "were", "what", "whatever", "when", "whence", "whenever", "where", "whereafter", "whereas", "whereby", "wherein", "whereupon", "wherever", "whether", "which", "while", "whither", "who", "whoever", "whole", "whom", "whose", "why", "will", "with", "within", "without", "would", "yet", "you", "your", "yours", "yourself", "yourselves"}
        words = re.findall(r"\b[a-zA-Z]{5,15}\b", text.lower())
        candidates = [w for w in words if w not in stopwords]
        freq = {}
        for c in candidates:
            freq[c] = freq.get(c, 0) + 1
        sorted_kws = sorted(freq.keys(), key=lambda x: freq[x], reverse=True)
        return sorted_kws[:5]

    @classmethod
    def _parse_txt(cls, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    @classmethod
    def _parse_markdown(cls, file_path: str) -> str:
        # Markdown is processed as plain text, chunks will parse headers nicely
        return cls._parse_txt(file_path)

    @classmethod
    def _parse_csv(cls, file_path: str) -> str:
        content = []
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for idx, row in enumerate(reader):
                content.append(", ".join(row))
        return "\n".join(content)

    @classmethod
    def _parse_json(cls, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            data = json.load(f)
            return json.dumps(data, indent=2)

    @classmethod
    def _parse_xml(cls, file_path: str) -> str:
        # Simple tag stripper for XML
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            raw_xml = f.read()
            clean_xml = re.sub(r"<[^>]+>", " ", raw_xml)
            return re.sub(r"\s+", " ", clean_xml).strip()

    @classmethod
    def _parse_html(cls, file_path: str) -> str:
        from bs4 import BeautifulSoup

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            html_content = f.read()

        soup = BeautifulSoup(html_content, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()

        clean_text = soup.get_text(separator=" ")
        return re.sub(r"\s+", " ", clean_text).strip()

    @classmethod
    def _parse_excel(cls, file_path: str) -> str:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text_runs = []
            for sheet in wb.worksheets:
                text_runs.append(f"--- Sheet: {sheet.title} ---")
                for row in sheet.iter_rows(values_only=True):
                    row_vals = [str(v) if v is not None else "" for v in row]
                    if any(row_vals):
                        text_runs.append(" | ".join(row_vals))
            return "\n".join(text_runs)
        except ImportError:
            logger.warning("openpyxl is not installed. Falling back to plain text extraction.")
            return cls._parse_txt(file_path)

    @classmethod
    def _parse_docx(cls, file_path: str) -> str:
        try:
            import docx
            doc = docx.Document(file_path)
            paragraphs = []
            for p in doc.paragraphs:
                if not p.text:
                    continue
                style_name = p.style.name if p.style else ""
                if style_name.startswith("Heading"):
                    lvl = style_name.replace("Heading", "").strip()
                    prefix = "#" * int(lvl) if lvl.isdigit() else "###"
                    paragraphs.append(f"{prefix} {p.text}")
                else:
                    paragraphs.append(p.text)
                    
            for table in doc.tables:
                table_rows = []
                for row_idx, row in enumerate(table.rows):
                    row_txt = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                    table_rows.append("| " + " | ".join(row_txt) + " |")
                    if row_idx == 0:
                        table_rows.append("| " + " | ".join(["---"] * len(row_txt)) + " |")
                if table_rows:
                    paragraphs.append("\n" + "\n".join(table_rows) + "\n")
                    
            return "\n".join(paragraphs)
        except ImportError:
            logger.warning("python-docx is not installed. Falling back to zipfile extraction.")
            try:
                import zipfile
                with zipfile.ZipFile(file_path) as z:
                    xml_content = z.read("word/document.xml").decode("utf-8")
                    clean_text = re.sub(r"<[^>]+>", " ", xml_content)
                    return re.sub(r"\s+", " ", clean_text).strip()
            except Exception:
                return cls._parse_txt(file_path)

    @classmethod
    def _parse_zip(cls, file_path: str, db: Session, organization_id: uuid.UUID, user_id: uuid.UUID) -> str:
        import zipfile
        import tempfile
        extracted_texts = []
        try:
            with zipfile.ZipFile(file_path, "r") as zip_ref:
                with tempfile.TemporaryDirectory() as tmp_dir:
                    zip_ref.extractall(tmp_dir)
                    for root, _, files in os.walk(tmp_dir):
                        for filename in files:
                            sub_path = os.path.join(root, filename)
                            sub_ext = filename.split(".")[-1].lower() if "." in filename else ""
                            if sub_ext in ["txt", "md", "csv", "json", "docx", "pdf", "html"]:
                                try:
                                    sub_res = cls.parse_file(sub_path, sub_ext, db, organization_id, user_id)
                                    extracted_texts.append(f"--- Archive File: {filename} ---\n{sub_res['content']}")
                                except Exception as e:
                                    logger.warning(f"Could not parse archive file {filename}: {e}")
            return "\n\n".join(extracted_texts) if extracted_texts else "ZIP archive contains no readable text files."
        except Exception as e:
            logger.error(f"Failed parsing ZIP archive {file_path}: {e}")
            return f"[ZIP archive error: {e}]"

    @classmethod
    def _parse_pptx(cls, file_path: str) -> str:
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            slides_text = []
            for idx, slide in enumerate(prs.slides):
                slides_text.append(f"--- Slide {idx + 1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slides_text.append(shape.text)
            return "\n".join(slides_text)
        except ImportError:
            logger.warning("python-pptx is not installed. Return empty content placeholder.")
            return f"[PowerPoint presentation: text extraction requires python-pptx]"

    @classmethod
    def _parse_rtf(cls, file_path: str) -> str:
        try:
            from striprtf.striprtf import rtf_to_text
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return rtf_to_text(f.read())
        except ImportError:
            # Basic regex tag stripper fallback for RTF
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                rtf = f.read()
                pattern = re.compile(r"\\([a-z]{1,32})(-?\d{1,10})?[ ]?|\\\'[0-9a-f]{2}|\\\{|\\\}|\\\n|[^\\\{\}\s]+")
                words = []
                for match in pattern.finditer(rtf):
                    word = match.group(0)
                    if not word.startswith("\\") and not word.startswith("{") and not word.startswith("}"):
                        words.append(word)
                return " ".join(words)

    @classmethod
    def _parse_image_ocr(
        cls,
        file_path: str,
        db: Session,
        organization_id: uuid.UUID,
        user_id: uuid.UUID,
    ) -> str:
        """
        Uses AI Gateway Vision model for image OCR.
        """
        try:
            from api.ai.gateway.coordinator import AIGateway
            gateway = AIGateway()
            
            with open(file_path, "rb") as image_file:
                encoded_string = base64.b64encode(image_file.read()).decode("utf-8")
                
            # Deduce extension mime type
            ext = os.path.splitext(file_path)[1].lower().strip(".")
            mime = f"image/{ext}" if ext in ["png", "jpg", "jpeg", "webp"] else "image/png"
            image_url = f"data:{mime};base64,{encoded_string}"
            
            prompt = (
                "Please perform OCR on this image. Extract all text content verbatim. "
                "Identify any tables, forms, figures, and handwritten text. "
                "Structure the output cleanly and preserve the reading order."
            )
            
            # Request vision inference via Gateway
            res = gateway.vision(
                db=db,
                prompt=prompt,
                image_url=image_url,
                organization_id=organization_id,
                user_id=user_id,
            )
            return res.get("content", "[OCR returned empty content]")
        except Exception as e:
            logger.error(f"Image OCR failed: {e}")
            raise RuntimeError(f"Image OCR processing error: {str(e)}")

    @classmethod
    def _parse_pdf(
        cls,
        file_path: str,
        db: Session,
        organization_id: uuid.UUID,
        user_id: uuid.UUID,
    ) -> str:
        """
        Check if PDF has text. If not, treat as Scanned PDF and run OCR using AI Gateway.
        """
        text_content = []
        has_text = False
        
        try:
            import pypdf
            reader = pypdf.PdfReader(file_path)
            for page in reader.pages:
                txt = page.extract_text()
                if txt and txt.strip():
                    text_content.append(txt)
                    has_text = True
        except ImportError:
            logger.warning("pypdf not installed. Falling back to native scan check.")
            
        if has_text:
            return "\n\n--- Page Break ---\n\n".join(text_content)
            
        # PDF is scanned (no extractable text) - run OCR using AI Gateway Vision.
        try:
            import pdf2image
            images = pdf2image.convert_from_path(file_path, first_page=1, last_page=5)
            ocr_pages = []
            
            from api.ai.gateway.coordinator import AIGateway
            gateway = AIGateway()
            
            for idx, img in enumerate(images):
                import io
                buffered = io.BytesIO()
                img.save(buffered, format="PNG")
                encoded_string = base64.b64encode(buffered.getvalue()).decode("utf-8")
                image_url = f"data:image/png;base64,{encoded_string}"
                
                prompt = f"Please transcribe Page {idx+1} of this scanned PDF document, extracting tables and text verbatim."
                res = gateway.vision(
                    db=db,
                    prompt=prompt,
                    image_url=image_url,
                    organization_id=organization_id,
                    user_id=user_id,
                )
                ocr_pages.append(res.get("content", ""))
                
            return "\n\n--- Page Break ---\n\n".join(ocr_pages)
        except Exception as e:
            logger.error(f"Scanned PDF OCR failed: {e}")
            raise RuntimeError(f"Scanned PDF processing failed: {str(e)}")
